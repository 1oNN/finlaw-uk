"""Document ingestion: parse uploaded files into text chunks and feed the
sparse retrieval index.

Supports: PDF (via pdfminer), DOCX (via python-docx), TXT, XLS/XLSX, PPTX.
"""

from __future__ import annotations

import os
from typing import List

from pdfminer.high_level import extract_text
import docx
import pandas as pd
from pptx import Presentation

from backend.retrieval.sparse import add_documents


def parse_pdf(path: str) -> List[str]:
    text = extract_text(path)
    return text.split("\f")


def parse_docx(path: str) -> List[str]:
    doc = docx.Document(path)
    return [p.text for p in doc.paragraphs]


def parse_txt(path: str) -> List[str]:
    with open(path, encoding="utf-8", errors="ignore") as f:
        return [f.read()]


def parse_excel(path: str) -> List[str]:
    df = pd.read_excel(path)
    return df.astype(str).values.flatten().tolist()


def parse_pptx(path: str) -> List[str]:
    prs = Presentation(path)
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    return texts


def ingest_file(path: str) -> int:
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        docs = parse_pdf(path)
    elif ext == ".docx":
        docs = parse_docx(path)
    elif ext == ".txt":
        docs = parse_txt(path)
    elif ext in (".xls", ".xlsx"):
        docs = parse_excel(path)
    elif ext == ".pptx":
        docs = parse_pptx(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")

    add_documents(docs)
    return len(docs)
